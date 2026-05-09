"""
Slide utilities: python-pptx helpers + Nano Banana Pro schematic generation.

This module is meant to be COPIED into your project and used by the Claude agent
to build slide decks. The agent imports this module and writes the main slide
script itself, tailored to the specific project content.

Requirements:
    pip install python-pptx google-genai Pillow

Environment:
    GOOGLE_API_KEY=your-gemini-api-key   (for Nano Banana Pro)
"""

import io
import os
import textwrap
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# =========================================================================
# SLIDE DIMENSIONS & PALETTE
# =========================================================================

SLIDE_WIDTH = Inches(13.333)   # Widescreen 16:9
SLIDE_HEIGHT = Inches(7.5)

# Clean modern palette
PALETTE = {
    'bg_dark':      RGBColor(0x1A, 0x1A, 0x2E),   # Deep navy
    'bg_mid':       RGBColor(0x16, 0x21, 0x3E),   # Dark blue
    'bg_light':     RGBColor(0xF8, 0xF9, 0xFA),   # Near white
    'accent':       RGBColor(0x4E, 0xC9, 0xB0),   # Teal
    'accent2':      RGBColor(0x56, 0x9C, 0xD6),   # Blue
    'accent3':      RGBColor(0xDC, 0xDC, 0xAA),   # Gold
    'text_light':   RGBColor(0xFF, 0xFF, 0xFF),   # White
    'text_dark':    RGBColor(0x22, 0x22, 0x22),   # Near black
    'text_muted':   RGBColor(0x99, 0x99, 0x99),   # Gray
    'highlight':    RGBColor(0xFF, 0x6B, 0x6B),   # Coral red
}


def create_presentation():
    """Create a blank widescreen 16:9 presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def _add_bg(slide, color):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 color=None, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri', line_spacing=1.2):
    """Add a text box to a slide with styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color or PALETTE['text_light']
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(font_size * 0.3)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def _add_multiline_textbox(slide, left, top, width, height, lines,
                            font_size=16, color=None, font_name='Calibri',
                            bullet=False, line_spacing=1.3):
    """Add a text box with multiple paragraphs/bullets."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if bullet:
            p.text = f"  {line}"
            p.level = 0
        else:
            p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color or PALETTE['text_light']
        p.font.name = font_name
        p.space_after = Pt(font_size * 0.4)
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def _add_accent_bar(slide, left, top, width, height, color=None):
    """Add a thin accent rectangle (decorative bar)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color or PALETTE['accent']
    shape.line.fill.background()
    return shape


# =========================================================================
# SLIDE TEMPLATES
# =========================================================================

def add_title_slide(prs, title, subtitle='', author=''):
    """Dark background title slide with accent bar."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _add_bg(slide, PALETTE['bg_dark'])

    # Accent bar
    _add_accent_bar(slide, Inches(0), Inches(3.0), Inches(13.333), Inches(0.06))

    # Title
    _add_textbox(slide, Inches(1.0), Inches(1.5), Inches(11.3), Inches(1.8),
                 title, font_size=44, bold=True, color=PALETTE['text_light'],
                 alignment=PP_ALIGN.LEFT)

    # Subtitle
    if subtitle:
        _add_textbox(slide, Inches(1.0), Inches(3.3), Inches(11.3), Inches(1.0),
                     subtitle, font_size=22, color=PALETTE['accent'],
                     alignment=PP_ALIGN.LEFT)

    # Author / date
    if author:
        _add_textbox(slide, Inches(1.0), Inches(5.8), Inches(11.3), Inches(0.6),
                     author, font_size=14, color=PALETTE['text_muted'],
                     alignment=PP_ALIGN.LEFT)
    return slide


def add_section_slide(prs, section_number, section_title, section_subtitle=''):
    """Dark section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, PALETTE['bg_mid'])

    # Big section number
    _add_textbox(slide, Inches(1.0), Inches(1.5), Inches(2.0), Inches(2.0),
                 f'{section_number:02d}', font_size=72, bold=True,
                 color=PALETTE['accent'])

    # Accent bar
    _add_accent_bar(slide, Inches(1.0), Inches(3.5), Inches(3.0), Inches(0.05))

    # Title
    _add_textbox(slide, Inches(1.0), Inches(3.8), Inches(11.0), Inches(1.5),
                 section_title, font_size=36, bold=True, color=PALETTE['text_light'])

    if section_subtitle:
        _add_textbox(slide, Inches(1.0), Inches(5.0), Inches(11.0), Inches(1.0),
                     section_subtitle, font_size=18, color=PALETTE['text_muted'])
    return slide


def add_content_slide(prs, title, bullets, footnote='', dark=True):
    """Standard content slide with title + bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, PALETTE['bg_dark'] if dark else PALETTE['bg_light'])
    text_color = PALETTE['text_light'] if dark else PALETTE['text_dark']

    # Title
    _add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.9),
                 title, font_size=28, bold=True, color=text_color)

    # Accent underline
    _add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2.0), Inches(0.04))

    # Bullets
    _add_multiline_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.0),
                            bullets, font_size=18, color=text_color, bullet=True)

    # Footnote
    if footnote:
        _add_textbox(slide, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.5),
                     footnote, font_size=11, color=PALETTE['text_muted'])
    return slide


def add_two_column_slide(prs, title, left_title, left_bullets,
                          right_title, right_bullets, dark=True):
    """Two-column content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, PALETTE['bg_dark'] if dark else PALETTE['bg_light'])
    text_color = PALETTE['text_light'] if dark else PALETTE['text_dark']

    # Title
    _add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.9),
                 title, font_size=28, bold=True, color=text_color)
    _add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2.0), Inches(0.04))

    # Left column
    _add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.6),
                 left_title, font_size=20, bold=True, color=PALETTE['accent'])
    _add_multiline_textbox(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(4.5),
                            left_bullets, font_size=16, color=text_color, bullet=True)

    # Divider
    _add_accent_bar(slide, Inches(6.6), Inches(1.5), Inches(0.03), Inches(5.0),
                     color=PALETTE['text_muted'])

    # Right column
    _add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.6),
                 right_title, font_size=20, bold=True, color=PALETTE['accent2'])
    _add_multiline_textbox(slide, Inches(7.0), Inches(2.1), Inches(5.5), Inches(4.5),
                            right_bullets, font_size=16, color=text_color, bullet=True)
    return slide


def add_key_metrics_slide(prs, title, metrics, dark=True):
    """Slide with big numbers / key metrics in a row.

    metrics: list of (value, label) tuples, e.g. [("95%", "Accuracy"), ("3.2s", "Latency")]
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, PALETTE['bg_dark'] if dark else PALETTE['bg_light'])
    text_color = PALETTE['text_light'] if dark else PALETTE['text_dark']

    _add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.9),
                 title, font_size=28, bold=True, color=text_color)
    _add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2.0), Inches(0.04))

    n = len(metrics)
    col_width = 11.7 / max(n, 1)
    for i, (value, label) in enumerate(metrics):
        x = Inches(0.8 + i * col_width)
        # Big number
        _add_textbox(slide, x, Inches(2.5), Inches(col_width), Inches(1.5),
                     str(value), font_size=52, bold=True, color=PALETTE['accent'],
                     alignment=PP_ALIGN.CENTER)
        # Label
        _add_textbox(slide, x, Inches(4.2), Inches(col_width), Inches(0.8),
                     label, font_size=16, color=PALETTE['text_muted'],
                     alignment=PP_ALIGN.CENTER)
    return slide


def add_image_slide(prs, image_path, title='', caption=''):
    """Full-slide image with optional small title overlay at top.

    For Nano Banana Pro schematics — the image IS the slide.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add image covering the full slide
    img_path = str(image_path)
    slide.shapes.add_picture(img_path, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT)

    # Optional semi-transparent title bar at top
    if title:
        # Dark overlay bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Emu(0), Emu(0), SLIDE_WIDTH, Inches(0.9))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
        # python-pptx doesn't support transparency directly on shape fill,
        # so we use a dark bar. For transparency, agent can post-process.
        bar.line.fill.background()

        _add_textbox(slide, Inches(0.5), Inches(0.1), Inches(12.3), Inches(0.7),
                     title, font_size=22, bold=True, color=PALETTE['text_light'])

    if caption:
        _add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.6),
                     caption, font_size=12, color=PALETTE['text_muted'])
    return slide


def add_closing_slide(prs, title='Thank You', contact='', message=''):
    """Closing/thank-you slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, PALETTE['bg_dark'])

    _add_accent_bar(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06))

    _add_textbox(slide, Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.5),
                 title, font_size=44, bold=True, color=PALETTE['text_light'],
                 alignment=PP_ALIGN.CENTER)

    if message:
        _add_textbox(slide, Inches(1.0), Inches(3.5), Inches(11.3), Inches(1.0),
                     message, font_size=20, color=PALETTE['accent'],
                     alignment=PP_ALIGN.CENTER)

    if contact:
        _add_textbox(slide, Inches(1.0), Inches(5.5), Inches(11.3), Inches(0.6),
                     contact, font_size=14, color=PALETTE['text_muted'],
                     alignment=PP_ALIGN.CENTER)
    return slide


# =========================================================================
# NANO BANANA PRO — SCHEMATIC GENERATION
# =========================================================================

def generate_schematic(prompt, output_path, aspect_ratio='16:9', size='2K'):
    """Generate a full-slide schematic image using Nano Banana Pro.

    Args:
        prompt: Detailed text prompt describing the schematic.
        output_path: Where to save the PNG.
        aspect_ratio: Image aspect ratio (default 16:9 for slides).
        size: Resolution — '1K', '2K', or '4K'.

    Returns:
        Path to the saved image.

    Requires:
        GOOGLE_API_KEY environment variable set.
    """
    from google import genai
    from google.genai import types

    api_key = os.environ.get('GOOGLE_API_KEY')
    if not api_key:
        raise RuntimeError(
            "GOOGLE_API_KEY not set. Export it before running:\n"
            "  export GOOGLE_API_KEY=your-key-here"
        )

    client = genai.Client(api_key=api_key)

    # Wrap user prompt with schematic-optimized instructions
    full_prompt = (
        "Create a clean, professional technical schematic diagram. "
        "Style: flat design with a dark background (#1a1a2e), using teal (#4ec9b0), "
        "blue (#569cd6), gold (#dcdcaa), and white for elements. "
        "The diagram should be visually clear with labeled components, arrows showing "
        "data/control flow, and a logical spatial layout. No photograph-like elements — "
        "this is an ABSTRACT TECHNICAL DIAGRAM suitable for a presentation slide. "
        "Minimize text in the image — use icons, shapes, and arrows instead. "
        f"\n\nSubject: {prompt}"
    )

    print(f"  Generating schematic: {Path(output_path).name}...")

    response = client.models.generate_content(
        model="gemini-3-pro-image-preview",
        contents=full_prompt,
        config=types.GenerateContentConfig(
            response_modalities=['IMAGE'],
            image_config=types.ImageConfig(
                aspect_ratio=aspect_ratio,
                image_size=size,
            ),
        ),
    )

    # Save the generated image
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    for part in response.parts:
        if part.inline_data is not None:
            image = part.as_image()
            image.save(str(output_path))
            print(f"  Saved: {output_path}")
            return str(output_path)

    raise RuntimeError(f"No image returned from Nano Banana Pro for prompt: {prompt[:80]}...")


def generate_schematics_batch(prompts_and_paths, aspect_ratio='16:9', size='2K'):
    """Generate multiple schematics. Returns list of saved paths.

    Args:
        prompts_and_paths: list of (prompt, output_path) tuples.
        aspect_ratio: Image aspect ratio.
        size: Resolution.

    Returns:
        List of paths to saved images.
    """
    results = []
    for i, (prompt, path) in enumerate(prompts_and_paths):
        print(f"  [{i+1}/{len(prompts_and_paths)}] Generating schematic...")
        try:
            saved = generate_schematic(prompt, path, aspect_ratio, size)
            results.append(saved)
        except Exception as e:
            print(f"  WARNING: Failed to generate schematic {i+1}: {e}")
            results.append(None)
    return results
